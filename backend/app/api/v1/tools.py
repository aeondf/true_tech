import io
import asyncio
from functools import partial

from fastapi import APIRouter, HTTPException, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, HttpUrl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from app.services.web_search import WebSearchService
from app.services.web_parser import WebParserService
from app.services.mws_client import MWSClient, get_mws_client
from app.config import Settings, get_settings

router = APIRouter()
_search = WebSearchService()
_parser = WebParserService()


class SearchRequest(BaseModel):
    query: str
    max_results: int = 5


class ParseRequest(BaseModel):
    url: HttpUrl
    extract_links: bool = False


class PptxRequest(BaseModel):
    topic: str
    slide_count: int = 7          # сколько слайдов генерировать
    language: str = "ru"           # язык презентации


@router.post("/web-search")
async def web_search(req: SearchRequest):
    """DuckDuckGo search → list of {title, url, snippet}."""
    results = await _search.search(req.query, max_results=req.max_results)
    return {"results": results}


@router.post("/web-parse")
async def web_parse(req: ParseRequest):
    """Fetch URL → extract clean text (+ optional links)."""
    result = await _parser.parse(str(req.url), extract_links=req.extract_links)
    return result


@router.post("/generate-pptx")
async def generate_pptx(
    req: PptxRequest,
    mws: MWSClient = Depends(get_mws_client),
    settings: Settings = Depends(get_settings),
):
    """
    Генерирует PowerPoint-презентацию по теме.
    1. LLM создаёт структуру слайдов в JSON.
    2. python-pptx собирает .pptx файл.
    3. Возвращает файл для скачивания.
    """
    # 1. Просим LLM сгенерировать структуру
    system_prompt = (
        "Ты генерируешь структуру презентации. "
        f"Создай ровно {req.slide_count} слайдов на языке: {req.language}. "
        "Ответь ТОЛЬКО валидным JSON-массивом без пояснений. "
        "Формат каждого слайда:\n"
        '{"title": "Заголовок слайда", "bullets": ["пункт 1", "пункт 2", "пункт 3"]}\n'
        "Первый слайд — титульный (bullets можно оставить пустым []).\n"
        "Последний слайд — итоги/выводы."
    )

    try:
        raw_json = await mws.chat_simple(
            model=settings.MODEL_TEXT,
            system=system_prompt,
            user=f"Тема презентации: {req.topic}",
        )
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"LLM недоступен: {e}")

    # Парсим JSON (LLM иногда оборачивает в ```json ... ```)
    import json, re
    json_match = re.search(r"\[.*\]", raw_json, re.DOTALL)
    if not json_match:
        raise HTTPException(status_code=500, detail="LLM вернул неверный формат JSON")

    try:
        slides_data: list[dict] = json.loads(json_match.group())
    except json.JSONDecodeError:
        raise HTTPException(status_code=500, detail="Не удалось распарсить структуру слайдов")

    # 2. Собираем .pptx в отдельном потоке (python-pptx синхронный)
    loop = asyncio.get_event_loop()
    pptx_bytes = await loop.run_in_executor(
        None, partial(_build_pptx, req.topic, slides_data)
    )

    # 3. Возвращаем файл
    filename = f"{req.topic[:40].replace(' ', '_')}.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def _build_pptx(topic: str, slides_data: list[dict]) -> bytes:
    """Синхронная сборка .pptx — запускается в thread pool."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Цветовая схема
    BG_COLOR = RGBColor(0x1E, 0x1E, 0x2E)      # тёмно-синий фон
    TITLE_COLOR = RGBColor(0xCB, 0xD3, 0xFF)    # светло-голубой заголовок
    TEXT_COLOR = RGBColor(0xE0, 0xE0, 0xFF)     # белый текст
    ACCENT_COLOR = RGBColor(0x89, 0xB4, 0xFA)   # акцент для буллетов

    for i, slide_info in enumerate(slides_data):
        layout = prs.slide_layouts[6]  # пустой layout
        slide = prs.slides.add_slide(layout)

        # Фон
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BG_COLOR

        slide_title = slide_info.get("title", f"Слайд {i + 1}")
        bullets = slide_info.get("bullets", [])

        # Заголовок
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(36 if i == 0 else 28)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR

        if not bullets:
            continue

        # Буллеты
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.9), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True

        for j, bullet_text in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {bullet_text}"
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR if j % 2 == 0 else ACCENT_COLOR
            p.space_after = Pt(8)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
