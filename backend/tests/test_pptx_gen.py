"""
Полный E2E тест генерации PPTX-презентаций:
  1. Генерация через backend API — POST /v1/tools/generate-pptx
  2. Проверка структуры файла (python-pptx)
  3. Разное кол-во слайдов
  4. Разные языки
  5. Проверка валидности PPTX

Запуск:
  python backend/tests/test_pptx_gen.py
"""
import time
import sys
import io
import httpx
import asyncio

BASE = "http://localhost:8000"
TIMEOUT = 120


def ok(name: str, status: int, elapsed: float, extra: str = ""):
    sym = "✅" if 200 <= status < 400 else "❌"
    print(f"  {sym} [{status}] {name} ({elapsed:.1f}s){' — ' + extra if extra else ''}")


def fail(name: str, msg: str):
    print(f"  ❌ {name} — {msg}")


async def run_tests():
    results = []

    # ── 1. Базовая генерация (7 слайдов, RU) ───────────────────
    print()
    print("═" * 60)
    print("1. Базовая генерация — 7 слайдов, русский")
    print("═" * 60)

    try:
        t0 = time.time()
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.post(f"{BASE}/v1/tools/generate-pptx", json={
                "topic": "Искусственный интеллект в медицине",
                "slide_count": 7,
                "language": "ru",
            })
        elapsed = time.time() - t0

        if r.status_code == 200:
            pptx_bytes = r.content
            ok("Basic PPTX (7 slides, RU)", r.status_code, elapsed, f"{len(pptx_bytes)} байт")

            # Проверяем структуру
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(pptx_bytes))
                n_slides = len(prs.slides)
                titles = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            titles.append(shape.text_frame.paragraphs[0].text[:50])
                            break

                print(f"    📊 Слайдов: {n_slides}")
                for i, t in enumerate(titles):
                    print(f"       {i+1}. {t}")

                # Сохраняем
                with open("/tmp/test_pptx_basic.pptx", "wb") as f:
                    f.write(pptx_bytes)
                print(f"    📁 Сохранён: /tmp/test_pptx_basic.pptx")
                results.append(("Basic PPTX (7 slides, RU)", True))
            except ImportError:
                print("    ⚠️  python-pptx не установлен для проверки, но файл получен")
                results.append(("Basic PPTX (7 slides, RU)", True))
        else:
            ok("Basic PPTX (7 slides, RU)", r.status_code, elapsed, r.text[:150])
            results.append(("Basic PPTX (7 slides, RU)", False))
    except Exception as e:
        fail("Basic PPTX (7 slides, RU)", str(e))
        results.append(("Basic PPTX (7 slides, RU)", False))

    # ── 2. Короткая презентация (3 слайда) ──────────────────────
    print()
    print("═" * 60)
    print("2. Короткая презентация — 3 слайда")
    print("═" * 60)

    try:
        t0 = time.time()
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.post(f"{BASE}/v1/tools/generate-pptx", json={
                "topic": "Квантовые компьютеры",
                "slide_count": 3,
                "language": "ru",
            })
        elapsed = time.time() - t0

        if r.status_code == 200:
            ok("Short PPTX (3 slides)", r.status_code, elapsed, f"{len(r.content)} байт")
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(r.content))
                print(f"    📊 Слайдов: {len(prs.slides)}")
            except ImportError:
                pass
            with open("/tmp/test_pptx_short.pptx", "wb") as f:
                f.write(r.content)
            print(f"    📁 Сохранён: /tmp/test_pptx_short.pptx")
            results.append(("Short PPTX (3 slides)", True))
        else:
            ok("Short PPTX (3 slides)", r.status_code, elapsed, r.text[:150])
            results.append(("Short PPTX (3 slides)", False))
    except Exception as e:
        fail("Short PPTX (3 slides)", str(e))
        results.append(("Short PPTX (3 slides)", False))

    # ── 3. Английский язык ──────────────────────────────────────
    print()
    print("═" * 60)
    print("3. English presentation — 5 slides")
    print("═" * 60)

    try:
        t0 = time.time()
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.post(f"{BASE}/v1/tools/generate-pptx", json={
                "topic": "Machine Learning in Finance",
                "slide_count": 5,
                "language": "en",
            })
        elapsed = time.time() - t0

        if r.status_code == 200:
            ok("English PPTX (5 slides)", r.status_code, elapsed, f"{len(r.content)} байт")
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(r.content))
                print(f"    📊 Slides: {len(prs.slides)}")
                for i, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            print(f"       {i+1}. {shape.text_frame.paragraphs[0].text[:60]}")
                            break
            except ImportError:
                pass
            with open("/tmp/test_pptx_en.pptx", "wb") as f:
                f.write(r.content)
            print(f"    📁 Saved: /tmp/test_pptx_en.pptx")
            results.append(("English PPTX (5 slides)", True))
        else:
            ok("English PPTX (5 slides)", r.status_code, elapsed, r.text[:150])
            results.append(("English PPTX (5 slides)", False))
    except Exception as e:
        fail("English PPTX (5 slides)", str(e))
        results.append(("English PPTX (5 slides)", False))

    # ── 4. Большая презентация (12 слайдов) ─────────────────────
    print()
    print("═" * 60)
    print("4. Большая презентация — 12 слайдов")
    print("═" * 60)

    try:
        t0 = time.time()
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.post(f"{BASE}/v1/tools/generate-pptx", json={
                "topic": "История развития интернета",
                "slide_count": 12,
                "language": "ru",
            })
        elapsed = time.time() - t0

        if r.status_code == 200:
            ok("Large PPTX (12 slides)", r.status_code, elapsed, f"{len(r.content)} байт")
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(r.content))
                print(f"    📊 Слайдов: {len(prs.slides)}")
            except ImportError:
                pass
            with open("/tmp/test_pptx_large.pptx", "wb") as f:
                f.write(r.content)
            print(f"    📁 Сохранён: /tmp/test_pptx_large.pptx")
            results.append(("Large PPTX (12 slides)", True))
        else:
            ok("Large PPTX (12 slides)", r.status_code, elapsed, r.text[:150])
            results.append(("Large PPTX (12 slides)", False))
    except Exception as e:
        fail("Large PPTX (12 slides)", str(e))
        results.append(("Large PPTX (12 slides)", False))

    # ── 5. Проверка Content-Disposition / MIME ──────────────────
    print()
    print("═" * 60)
    print("5. Проверка заголовков ответа")
    print("═" * 60)

    try:
        t0 = time.time()
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.post(f"{BASE}/v1/tools/generate-pptx", json={
                "topic": "Тест заголовков",
                "slide_count": 3,
            })
        elapsed = time.time() - t0

        if r.status_code == 200:
            ct = r.headers.get("content-type", "")
            cd = r.headers.get("content-disposition", "")
            mime_ok = "officedocument.presentationml.presentation" in ct
            disp_ok = "attachment" in cd and ".pptx" in cd
            ok("Response headers", r.status_code, elapsed,
               f"MIME={'✅' if mime_ok else '❌'} Disposition={'✅' if disp_ok else '❌'}")
            print(f"    Content-Type: {ct}")
            print(f"    Content-Disposition: {cd}")
            results.append(("Response headers", mime_ok and disp_ok))
        else:
            ok("Response headers", r.status_code, elapsed, r.text[:100])
            results.append(("Response headers", False))
    except Exception as e:
        fail("Response headers", str(e))
        results.append(("Response headers", False))

    # ── Summary ─────────────────────────────────────────────────
    print()
    print("═" * 60)
    print("ИТОГО — PPTX Генерация")
    print("═" * 60)
    for name, passed in results:
        if passed:
            print(f"  ✅ {name}")
        else:
            print(f"  ❌ {name}")

    total = sum(1 for _, p in results if p)
    failed = sum(1 for _, p in results if not p)
    print(f"\n  Passed: {total}  Failed: {failed}")


if __name__ == "__main__":
    asyncio.run(run_tests())
